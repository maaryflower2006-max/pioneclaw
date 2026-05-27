"""
内置工具集 - 提供常用的基础工具

包含：
- EchoTool: 回显工具（测试用）
- CurrentTimeTool: 获取当前时间
- CalculatorTool: 简单计算器
"""

import ast
import asyncio
import glob
import json
import logging
import operator
import os
import re
import shlex
import uuid
from datetime import datetime
from typing import TYPE_CHECKING, Optional

from app.core.bash_safety import (
    CommandConfirmationRequired,
    DangerLevel,
    analyze_command,
)
from app.core.sandbox import (
    PathOutsideWorkspaceError,
    SensitiveFileAccessRequired,
)
from app.modules.tools.base import BaseTool, ToolParameter
from app.modules.tools.mcp import ListMcpResourcesTool, McpAuthTool, MCPTool, ReadMcpResourceTool
from app.modules.tools.plan_mode import EnterPlanModeTool, ExitPlanModeTool
from app.modules.tools.registry import get_tool_registry
from app.modules.tools.send_message import (
    register_agent,
    unregister_agent,
)
from app.modules.tools.task_store import (
    _background_tasks,
)
from app.modules.tools.task_store import (
    create_task as _store_create_task,
)
from app.modules.tools.task_store import (
    get_task as _store_get_task,
)
from app.modules.tools.task_store import (
    list_tasks as _store_list_tasks,
)
from app.modules.tools.web import WebFetchTool, WebSearchTool

if TYPE_CHECKING:
    from app.modules.tools.registry import ToolRegistry

logger = logging.getLogger(__name__)


class CurrentTimeTool(BaseTool):
    """获取当前时间"""

    name = "current_time"
    description = "获取当前日期和时间"
    is_parallel_safe = True
    parameters = {
        "timezone": ToolParameter(
            type="string",
            description="时区，如 'Asia/Shanghai'，默认为本地时区",
            default="local",
        ),
    }
    required = []

    async def execute(self, timezone: str = "local", **kwargs) -> str:
        now = datetime.now()
        return f"当前时间: {now.strftime('%Y-%m-%d %H:%M:%S')} ({timezone})"


# AST-based safe math evaluator (replaces eval())
_SAFE_BINOPS = {
    ast.Add: operator.add,
    ast.Sub: operator.sub,
    ast.Mult: operator.mul,
    ast.Div: operator.truediv,
    ast.FloorDiv: operator.floordiv,
    ast.Mod: operator.mod,
    ast.Pow: operator.pow,
    ast.BitXor: operator.pow,  # ^ 当作幂运算（如 5^2 = 25）
}

def _safe_eval_node(node):
    """递归安全求值 AST 节点，仅允许数学运算"""
    if isinstance(node, ast.Constant):
        if isinstance(node.value, (int, float)):
            return node.value
        raise ValueError(f"不允许的常量类型: {type(node.value).__name__}")
    if isinstance(node, ast.BinOp):
        left = _safe_eval_node(node.left)
        right = _safe_eval_node(node.right)
        op = _SAFE_BINOPS.get(type(node.op))
        if op is None:
            raise ValueError(f"不支持的运算符: {type(node.op).__name__}")
        return op(left, right)
    if isinstance(node, ast.UnaryOp):
        operand = _safe_eval_node(node.operand)
        if isinstance(node.op, ast.USub):
            return -operand
        if isinstance(node.op, ast.UAdd):
            return +operand
        raise ValueError(f"不支持的一元运算符: {type(node.op).__name__}")
    raise ValueError(f"不支持的表达式: {type(node).__name__}")

def safe_eval_math(expression: str):
    """仅允许数学表达式求值，无代码执行风险"""
    tree = ast.parse(expression.strip(), mode="eval")
    return _safe_eval_node(tree.body)


class CalculatorTool(BaseTool):
    """安全计算器 — 基于 AST 解析，无 eval() 风险"""

    name = "calculator"
    description = "执行数学计算，支持 + - * / // % ** ()，如 '2 + 3 * 4'"
    is_parallel_safe = True
    parameters = {
        "expression": ToolParameter(
            type="string",
            description="数学表达式，如 '2 + 3 * 4'",
        ),
    }
    required = ["expression"]

    async def execute(self, expression: str, **kwargs) -> str:
        try:
            result = safe_eval_math(expression)
            if isinstance(result, float) and result == int(result):
                result = int(result)
            return f"计算结果: {expression} = {result}"
        except ZeroDivisionError:
            return "错误: 除以零"
        except Exception as e:
            return f"计算错误: {e}"


class ThinkingTool(BaseTool):
    """显式思考工具 - 让 AI 输出推理过程"""

    name = "thinking"
    description = (
        "Use this tool to think through complex problems step by step. "
        "Before making significant changes (refactoring, architecture decisions, "
        "complex debugging), use this tool to plan your approach. "
        "The thought process will be visible to the user but won't affect the system."
    )
    parameters = {
        "thought": ToolParameter(
            type="string",
            description="Your step-by-step reasoning process. Be explicit about assumptions, trade-offs, and plan.",
        ),
    }
    required = ["thought"]
    is_parallel_safe = True

    async def execute(self, thought: str, **kwargs) -> str:
        preview = thought[:100].replace("\n", " ")
        if len(thought) > 100:
            preview += "..."
        return f"[Thinking] {len(thought)} characters | Preview: {preview}"


class ViewTool(BaseTool):
    """智能查看工具 - 自动判断文件或目录，支持渐进式阅读"""

    name = "view"
    is_parallel_safe = True
    description = (
        "View the contents of a file or directory. "
        "If the path is a file, returns its content (with optional offset/limit). "
        "If the path is a directory, lists its contents. "
        "This is the primary tool for exploring the codebase."
    )
    parameters = {
        "path": ToolParameter(
            type="string",
            description="File or directory path",
        ),
        "offset": ToolParameter(
            type="integer",
            description="Line offset for files (1-based), default 1",
            default=1,
        ),
        "limit": ToolParameter(
            type="integer",
            description="Max lines to read for files, default 50",
            default=50,
        ),
    }
    required = ["path"]

    async def execute(self, path: str, offset: int = 1, limit: int = 50, **kwargs) -> str:
        from pathlib import Path

        from app.core.config import settings
        from app.core.sandbox import validate_path_for_read

        safe_path = validate_path_for_read(
            path,
            Path(settings.WORKSPACE_DIR),
            allow_sensitive=kwargs.get("_allow_sensitive", False),
            allow_outside=kwargs.get("_allow_outside", False),
        )

        if safe_path.is_dir():
            return await self._view_directory(safe_path, limit)
        else:
            return await self._view_file(safe_path, offset, limit)

    async def _view_file(self, safe_path, offset: int, limit: int) -> str:
        rf = ReadFileTool()
        return await rf.execute(
            path=str(safe_path),
            offset=offset,
            limit=limit,
            raw=False,
        )

    async def _view_directory(self, path, limit: int) -> str:
        entries = []
        try:
            with os.scandir(str(path)) as it:
                for entry in it:
                    try:
                        if entry.is_dir():
                            entries.append(("[DIR] ", entry.name, ""))
                        else:
                            size = entry.stat().st_size
                            if size > 1024 * 1024:
                                size_str = f" ({size / (1024 * 1024):.1f}MB)"
                            elif size > 1024:
                                size_str = f" ({size / 1024:.0f}KB)"
                            else:
                                size_str = f" ({size}B)"
                            entries.append(("[FILE]", entry.name, size_str))
                    except OSError:
                        entries.append(("[???] ", entry.name, ""))
        except OSError as e:
            return f"无法读取目录: {e}"

        if not entries:
            return f"目录为空: {path}"

        entries.sort(key=lambda e: (e[0] != "[DIR] ", e[1].lower()))
        total = len(entries)
        if limit < total:
            entries = entries[:limit]

        lines = [f"{path} ({total} 个条目):"]
        for prefix, name, size in entries:
            lines.append(f"{prefix}{name}{size}")

        if limit < total:
            lines.append(f"... (仅显示前 {limit} 个)")

        return "\n".join(lines)


class ReadFileTool(BaseTool):
    """读取文件工具"""

    name = "read_file"
    description = (
        "读取文件内容，支持文本文件（.txt, .md, .csv, .json 等）、Word（.docx）、"
        "Excel（.xlsx, .xls）、PowerPoint（.pptx）和 PDF（.pdf）。"
        "支持 offset/limit 分页阅读。"
    )
    is_parallel_safe = True
    parameters = {
        "path": ToolParameter(
            type="string",
            description="文件路径",
        ),
        "offset": ToolParameter(
            type="integer",
            description="起始行号（1-based），默认 1",
            default=1,
        ),
        "limit": ToolParameter(
            type="integer",
            description="最大读取行数，默认 200",
            default=200,
        ),
        "sheet_name": ToolParameter(
            type="string",
            description="Excel 工作表名称（仅 Excel 文件有效），默认为第一个工作表",
            default="",
        ),
        "max_rows": ToolParameter(
            type="integer",
            description="最大读取行数（适用于 Excel/CSV），默认 200",
            default=200,
        ),
        "raw": ToolParameter(
            type="boolean",
            description="是否返回原始纯文本（不加标题行号前缀）。默认 true 保持向后兼容。",
            default=True,
        ),
    }
    required = ["path"]

    async def execute(self, path: str, *, offset: int = 1, limit: int = 200, sheet_name: str = "", max_rows: int = 200, raw: bool = True, **kwargs) -> str:
        try:
            # 沙箱校验
            from pathlib import Path

            from app.core.config import settings
            from app.core.sandbox import validate_path_for_read
            safe_path = validate_path_for_read(
                path, Path(settings.WORKSPACE_DIR),
            )

            ext = os.path.splitext(str(safe_path))[1].lower()

            sp = str(safe_path)
            # Excel 文件处理
            if ext in ('.xlsx', '.xls'):
                return await self._read_excel(sp, sheet_name, max_rows)

            # Word 文件处理
            if ext == '.docx':
                return await self._read_docx(sp)

            # PowerPoint 文件处理
            if ext == '.pptx':
                return await self._read_pptx(sp)

            # PDF 文件处理
            if ext == '.pdf':
                return await self._read_pdf(sp, max_rows)

            # CSV 文件处理（用 utf-8 读）
            if ext == '.csv':
                return await self._read_csv(sp, max_rows)

            # 其它文本文件（支持 offset/limit 分页）
            try:
                with open(sp, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                # 尝试其他编码
                with open(sp, 'r', encoding='gbk') as f:
                    content = f.read()

            lines = content.splitlines()
            total = len(lines)
            start = max(0, offset - 1)
            end = min(total, start + limit)
            selected = lines[start:end]

            if not raw:
                header = f"[{os.path.basename(sp)}] Lines {start + 1}-{end} of {total}:"
                selected.insert(0, header)
            if end < total:
                selected.append("... (内容已截断，使用 offset/limit 查看更多)")
            result = "\n".join(selected)
            if len(result) > 8000:
                result = result[:8000] + "\n... (内容已截断)"
            return result

        except FileNotFoundError:
            return f"错误: 文件不存在 - {path}"
        except (PathOutsideWorkspaceError, SensitiveFileAccessRequired):
            raise
        except Exception as e:
            return f"读取文件错误: {e}"

    async def _read_excel(self, path: str, sheet_name: str, max_rows: int) -> str:
        """读取 Excel 文件并转换为 Markdown 表格文本"""
        try:
            import openpyxl
        except ImportError:
            return "错误: 需要安装 openpyxl 库才能读取 Excel 文件 (pip install openpyxl)"

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)

        # 选择工作表
        if sheet_name:
            if sheet_name not in wb.sheetnames:
                return f"错误: 工作表 '{sheet_name}' 不存在。可用工作表: {', '.join(wb.sheetnames)}"
            ws = wb[sheet_name]
        else:
            ws = wb.active

        # 读取数据
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            rows.append([str(cell) if cell is not None else '' for cell in row])

        wb.close()

        if not rows:
            return "Excel 文件为空"

        # 生成文本表格
        lines = [f"工作表: {ws.title} ({len(rows)} 行, {len(rows[0])} 列)", ""]
        lines.append(" | ".join(rows[0]))
        lines.append(" | ".join(["---"] * len(rows[0])))
        for row in rows[1:]:
            # 确保每行列数一致
            padded = row + [''] * (len(rows[0]) - len(row))
            lines.append(" | ".join(padded[:len(rows[0])]))

        if len(rows) >= max_rows:
            lines.append(f"\n... (仅显示前 {max_rows} 行)")

        result = "\n".join(lines)
        if len(result) > 8000:
            result = result[:8000] + "\n... (内容已截断)"
        return result

    async def _read_csv(self, path: str, max_rows: int) -> str:
        """读取 CSV 文件"""
        lines = []
        with open(path, 'r', encoding='utf-8') as f:
            for i, line in enumerate(f):
                if i >= max_rows:
                    break
                lines.append(line.rstrip('\n'))

        if not lines:
            return "CSV 文件为空"

        result = "\n".join(lines)
        if len(lines) >= max_rows:
            result += f"\n... (仅显示前 {max_rows} 行)"
        if len(result) > 8000:
            result = result[:8000] + "\n... (内容已截断)"
        return result

    async def _read_docx(self, path: str) -> str:
        """读取 Word 文档（.docx）"""
        try:
            from docx import Document
        except ImportError:
            return "错误: 需要安装 python-docx 库才能读取 Word 文件 (pip install python-docx)"

        doc = Document(path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
            elif paragraphs and paragraphs[-1] != '':
                paragraphs.append('')  # 空段落作为分隔

        if not paragraphs:
            return "Word 文档为空"

        result = "\n".join(paragraphs)
        if len(result) > 8000:
            result = result[:8000] + "\n... (内容已截断)"
        return result

    async def _read_pptx(self, path: str) -> str:
        """读取 PowerPoint 演示文稿（.pptx）"""
        try:
            from pptx import Presentation
        except ImportError:
            return "错误: 需要安装 python-pptx 库才能读取 PowerPoint 文件 (pip install python-pptx)"

        prs = Presentation(path)
        slides_output = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_output.append(f"## 第 {i} 页\n" + "\n".join(texts))

        if not slides_output:
            return "PowerPoint 文件为空"

        result = "\n\n".join(slides_output)
        if len(result) > 8000:
            result = result[:8000] + "\n... (内容已截断)"
        return result

    async def _read_pdf(self, path: str, max_rows: int = 200) -> str:
        """读取 PDF 文件，提取所有页面文本"""
        try:
            import pdfplumber
        except ImportError:
            return "错误: 需要安装 pdfplumber 库才能读取 PDF 文件 (pip install pdfplumber)"

        with pdfplumber.open(path) as pdf:
            if len(pdf.pages) == 0:
                return "PDF 文件为空"

            pages_output = []
            total_lines = 0
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    lines = text.strip().split('\n')
                    if total_lines + len(lines) > max_rows:
                        lines = lines[:max_rows - total_lines]
                        pages_output.append(f"## 第 {i} 页\n" + "\n".join(lines))
                        pages_output.append("... (已达到最大行数限制)")
                        break
                    pages_output.append(f"## 第 {i} 页\n" + "\n".join(lines))
                    total_lines += len(lines)

        if not pages_output:
            return "PDF 文件无可提取文本（可能是扫描件或图片型 PDF）"

        result = "\n\n".join(pages_output)
        if len(result) > 8000:
            result = result[:8000] + "\n... (内容已截断)"
        return result


class WriteFileTool(BaseTool):
    """写入文件工具"""

    name = "write_file"
    description = "将内容写入文件。支持纯文本、Markdown 表格转 Excel（.xlsx）、Word（.docx）、PowerPoint（.pptx）"
    parameters = {
        "path": ToolParameter(
            type="string",
            description="文件路径",
        ),
        "content": ToolParameter(
            type="string",
            description="要写入的内容",
        ),
    }
    required = ["path", "content"]

    async def execute(self, path: str, content: str, **kwargs) -> str:
        try:
            # 沙箱校验
            from pathlib import Path

            from app.core.config import settings
            from app.core.sandbox import validate_path_for_write
            safe_path = validate_path_for_write(
                path, Path(settings.WORKSPACE_DIR),
                allow_sensitive=kwargs.get("_allow_sensitive", False),
            )

            sp = str(safe_path)
            ext = os.path.splitext(sp)[1].lower()

            if ext in ('.xlsx', '.xls'):
                return await self._write_xlsx(sp, content)
            if ext == '.docx':
                return await self._write_docx(sp, content)
            if ext == '.pptx':
                return await self._write_pptx(sp, content)

            # 纯文本写入
            with open(sp, 'w', encoding='utf-8') as f:
                f.write(content)
            return f"成功写入文件: {sp} ({len(content)} 字符)"
        except (PathOutsideWorkspaceError, SensitiveFileAccessRequired):
            raise
        except Exception as e:
            return f"写入文件错误: {e}"

    async def _write_xlsx(self, path: str, content: str) -> str:
        """将内容写入 Excel 文件"""
        try:
            import openpyxl
        except ImportError:
            return "错误: 需要安装 openpyxl 库才能写入 Excel 文件 (pip install openpyxl)"

        # 解析内容为行
        lines = [line for line in content.strip().split('\n') if line.strip()]
        rows = []

        # 检测是否为 markdown table
        if any('|' in line for line in lines):
            # Markdown table 模式：跳过分隔行（如 |---|）
            for line in lines:
                line = line.strip()
                if line.startswith('|') and '---' in line:
                    continue
                cells = [c.strip() for c in line.strip('|').split('|')]
                if any(c for c in cells):  # 跳过全空行
                    rows.append(cells)
        else:
            # CSV/TSV 模式
            for line in lines:
                if '\t' in line:
                    rows.append(line.split('\t'))
                elif ',' in line:
                    rows.append(line.split(','))
                else:
                    rows.append([line])

        if not rows:
            return "错误: 无法解析内容为表格数据"

        wb = openpyxl.Workbook()
        ws = wb.active
        for row in rows:
            ws.append(row)
        wb.save(path)

        return f"成功写入 Excel: {path} ({len(rows)} 行, {len(rows[0])} 列)"

    async def _write_docx(self, path: str, content: str) -> str:
        """将内容写入 Word 文档"""
        try:
            from docx import Document
        except ImportError:
            return "错误: 需要安装 python-docx 库才能写入 Word 文件 (pip install python-docx)"

        doc = Document()
        paragraphs = content.strip().split('\n\n')
        for para_text in paragraphs:
            para_text = para_text.strip()
            if para_text:
                doc.add_paragraph(para_text)
        doc.save(path)

        return f"成功写入 Word: {path} ({len(paragraphs)} 段)"

    async def _write_pptx(self, path: str, content: str) -> str:
        """将内容写入 PowerPoint 演示文稿"""
        try:
            from pptx import Presentation
        except ImportError:
            return "错误: 需要安装 python-pptx 库才能写入 PowerPoint 文件 (pip install python-pptx)"

        prs = Presentation()
        slides = content.strip().split('\n---\n')
        if len(slides) == 1:
            # 没有分隔符，尝试用空行分隔（标题独占一行模式）
            slides = [content.strip()]

        for slide_text in slides:
            slide_text = slide_text.strip()
            if not slide_text:
                continue
            lines = slide_text.split('\n')
            title = lines[0].strip()
            body = '\n'.join(line.strip() for line in lines[1:] if line.strip())

            slide_layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title if title else "(无标题)"
            if body and len(slide.placeholders) > 1:
                slide.placeholders[1].text = body

        prs.save(path)
        return f"成功写入 PowerPoint: {path} ({len(prs.slides)} 页)"


class ExecTool(BaseTool):
    """Shell 命令执行工具"""

    name = "exec"
    description = "在工作目录中执行 Shell 命令，返回输出结果。超时 30 秒，输出上限 5000 字符。危险命令（rm -rf、shutdown 等）被阻止或需要确认"
    parameters = {
        "command": ToolParameter(
            type="string",
            description="要执行的 Shell 命令",
        ),
    }
    required = ["command"]

    # Windows CMD 内置命令列表（create_subprocess_exec 无法直接执行这些）
    _WIN_CMD_BUILTINS = {
        "dir", "echo", "type", "cd", "chdir", "md", "mkdir", "rd", "rmdir",
        "del", "erase", "copy", "move", "ren", "rename", "cls", "pause",
        "ver", "vol", "path", "set", "exit", "start", "assoc", "ftype",
        "prompt", "title", "color", "date", "time",
    }

    async def execute(self, command: str, **kwargs) -> str:
        # 解析命令为参数列表（避免 shell 注入）
        # Windows 使用 posix=False 避免反斜杠被当作转义符，
        # 否则 C:\Users\file.txt 会被错误切分为 C:Usersfile.txt
        try:
            import sys
            is_windows = sys.platform == "win32"
            posix_mode = not is_windows
            args = shlex.split(command, posix=posix_mode)
            # 非 POSIX 模式会保留引号，这里手动去外层引号以保持行为一致
            if is_windows:
                args = [a[1:-1] if len(a) >= 2 and a[0] == a[-1] == '"' else a for a in args]
        except ValueError as e:
            return f"错误: 命令解析失败 - {e}"

        if not args:
            return "错误: 空命令"

        # 安全检查
        assessment = analyze_command(command)

        if assessment.level == DangerLevel.BLOCKED:
            return f"错误: 命令被安全策略拦截 - {assessment.risk_summary}"

        if assessment.level in (DangerLevel.CAUTION, DangerLevel.DANGEROUS):
            raise CommandConfirmationRequired(assessment)

        # 用 create_subprocess_exec 替代 create_subprocess_shell
        # shell=False 确保元字符（| ; && $() ``）不会被解释
        from pathlib import Path

        from app.core.config import settings
        workspace = Path(settings.WORKSPACE_DIR).resolve()
        workspace.mkdir(parents=True, exist_ok=True)
        cwd = str(workspace)

        import sys
        is_windows = sys.platform == "win32"

        # Windows 兼容性：CMD 内置命令（dir、echo、type 等）无法通过 create_subprocess_exec 直接执行
        if is_windows and args and args[0].lower() in self._WIN_CMD_BUILTINS:
            args = ["cmd", "/c", command]

        async def _run_subprocess(cmd_args):
            process = await asyncio.create_subprocess_exec(
                *cmd_args,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=cwd,
            )
            try:
                stdout, stderr = await asyncio.wait_for(
                    process.communicate(), timeout=30
                )
            except asyncio.TimeoutError:
                process.kill()
                await process.wait()
                raise asyncio.TimeoutError("命令执行超过 30 秒")

            output_parts = []
            if stdout:
                text = stdout.decode('utf-8', errors='replace')
                text = text.replace('\r\n', '\n').replace('\r', '\n')
                output_parts.append(text)
            if stderr:
                text = stderr.decode('utf-8', errors='replace')
                text = text.replace('\r\n', '\n').replace('\r', '\n')
                if text.strip():
                    output_parts.append(f"STDERR:\n{text}")
            if process.returncode != 0:
                output_parts.append(f"Exit code: {process.returncode}")

            result = "\n".join(output_parts) if output_parts else "(无输出)"
            if len(result) > 5000:
                result = result[:5000] + "\n... (输出已截断)"
            return result

        try:
            return await _run_subprocess(args)
        except FileNotFoundError:
            # Windows 回退：如果找不到可执行文件，尝试用 cmd /c 包装
            if is_windows and args[0:1] != ["cmd"]:
                try:
                    return await _run_subprocess(["cmd", "/c", command])
                except FileNotFoundError:
                    pass
            return f"错误: 找不到可执行文件 '{args[0]}'，请检查命令是否正确"
        except asyncio.TimeoutError:
            return "错误: 命令超时 (30秒)"
        except Exception as e:
            # 确保错误信息永远不会为空：同时输出异常类型和消息
            err_type = type(e).__name__
            err_msg = str(e) or "(无详细错误信息)"
            return f"执行命令错误 [{err_type}]: {err_msg}"


class EditFileTool(BaseTool):
    """编辑文件工具 — 精确文本替换"""

    name = "edit_file"
    description = (
        "编辑文件内容。提供 old_text 和 new_text，将文件中唯一的 old_text 精确替换为 new_text。"
        "old_text 必须在文件中唯一存在。编辑前建议先用 read_file 确认内容。"
    )
    parameters = {
        "path": ToolParameter(
            type="string",
            description="要编辑的文件路径",
        ),
        "old_text": ToolParameter(
            type="string",
            description="要替换的原始文本（文本替换模式）",
            default="",
        ),
        "new_text": ToolParameter(
            type="string",
            description="替换后的新文本",
            default="",
        ),
        "start_line": ToolParameter(
            type="integer",
            description="起始行号（1-based，行编辑模式）",
            default=0,
        ),
        "end_line": ToolParameter(
            type="integer",
            description="结束行号（默认等于 start_line）",
            default=0,
        ),
        "insert": ToolParameter(
            type="boolean",
            description="在 start_line 前插入而非替换（默认 false）",
            default=False,
        ),
    }
    required = ["path"]

    async def execute(self, path: str, old_text: str = "", new_text: str = "",
                      start_line: int = 0, end_line: int = 0,
                      insert: bool = False, **kwargs) -> str:
        try:
            # 沙箱校验（编辑 = 写操作）
            from pathlib import Path

            from app.core.config import settings
            from app.core.sandbox import validate_path_for_write
            safe_path = validate_path_for_write(
                path, Path(settings.WORKSPACE_DIR),
                allow_sensitive=kwargs.get("_allow_sensitive", False),
            )

            sp = str(safe_path)
            if not old_text:
                return "错误: 需要提供 old_text"
            return await self._edit_by_text(sp, old_text, new_text)

        except FileNotFoundError:
            return f"错误: 文件不存在 - {path}"
        except (PathOutsideWorkspaceError, SensitiveFileAccessRequired):
            raise
        except Exception as e:
            return f"编辑文件错误: {e}"

    async def _edit_by_text(self, path: str, old_text: str, new_text: str) -> str:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()

        if old_text not in content:
            total = len(content.splitlines())
            return (
                f"错误: 未找到 old_text（{total} 行，{len(content)} 字符）。"
                f"请用 read_file 确认文件精确内容"
            )

        count = content.count(old_text)
        if count > 1:
            return f"警告: old_text 出现 {count} 次，请增加上下文使其唯一"

        new_content = content.replace(old_text, new_text, 1)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(new_content)
        return f"已编辑 {path}（替换 1 处）"

class FileSearchTool(BaseTool):
    """文件搜索工具 — 按通配符模式搜索文件"""

    name = "file_search"
    description = "在指定目录中按通配符模式搜索文件（如 *.py、test_*.txt）"
    parameters = {
        "path": ToolParameter(
            type="string",
            description="搜索目录路径",
        ),
        "pattern": ToolParameter(
            type="string",
            description="文件名通配符模式（* 匹配任意字符，? 匹配单个字符），默认 *",
            default="*",
        ),
        "recursive": ToolParameter(
            type="boolean",
            description="是否递归搜索子目录，默认 true",
            default=True,
        ),
        "limit": ToolParameter(
            type="integer",
            description="最大返回结果数（1-100），默认 30",
            default=30,
        ),
    }
    required = ["path"]

    async def execute(self, path: str, pattern: str = "*", recursive: bool = True,
                      limit: int = 30, **kwargs) -> str:
        try:
            from pathlib import Path

            from app.core.config import settings
            from app.core.sandbox import validate_path_for_read
            safe_path = validate_path_for_read(
                path, Path(settings.WORKSPACE_DIR),
                allow_sensitive=kwargs.get("_allow_sensitive", False),
                allow_outside=kwargs.get("_allow_outside", False),
            )
            path = str(safe_path)
            if not os.path.isdir(path):
                return f"错误: 目录不存在 - {path}"

            if recursive:
                search_pattern = os.path.join(path, "**", pattern)
            else:
                search_pattern = os.path.join(path, pattern)

            results = []
            for f in glob.iglob(search_pattern, recursive=recursive):
                if len(results) >= limit:
                    break
                if os.path.isfile(f) or os.path.isdir(f):
                    ftype = "DIR" if os.path.isdir(f) else "FILE"
                    size = os.path.getsize(f) if os.path.isfile(f) else 0
                    rel = os.path.relpath(f, path)
                    if size > 1024 * 1024:
                        size_str = f"{size / (1024*1024):.1f}MB"
                    elif size > 1024:
                        size_str = f"{size / 1024:.0f}KB"
                    else:
                        size_str = f"{size}B"
                    results.append(f"[{ftype}] {rel} ({size_str})")

            if not results:
                return f"未找到匹配 '{pattern}' 的文件（{path}）"

            result = f"搜索 {path}，模式 '{pattern}'，找到 {len(results)} 个结果：\n" + "\n".join(results)
            if len(results) >= limit:
                result += f"\n... (结果已达上限 {limit})"
            return result

        except Exception as e:
            return f"搜索文件错误: {e}"


class ListDirTool(BaseTool):
    """列出目录内容"""

    name = "list_dir"
    is_parallel_safe = True
    description = "列出指定目录的内容，显示文件和子目录的名称、类型和大小"
    parameters = {
        "path": ToolParameter(
            type="string",
            description="要列出的目录路径",
        ),
        "limit": ToolParameter(
            type="integer",
            description="最大返回条目数（1-200），默认 50",
            default=50,
        ),
    }
    required = ["path"]

    async def execute(self, path: str, limit: int = 50, **kwargs) -> str:
        try:
            from pathlib import Path

            from app.core.config import settings
            from app.core.sandbox import validate_path_for_read
            safe_path = validate_path_for_read(
                path, Path(settings.WORKSPACE_DIR),
                allow_sensitive=kwargs.get("_allow_sensitive", False),
                allow_outside=kwargs.get("_allow_outside", False),
            )
            path = str(safe_path)
            if not os.path.isdir(path):
                return f"错误: 目录不存在 - {path}"

            entries = []
            with os.scandir(path) as it:
                for entry in it:
                    try:
                        if entry.is_dir():
                            entries.append(("[DIR] ", entry.name, ""))
                        else:
                            size = entry.stat().st_size
                            if size > 1024 * 1024:
                                size_str = f" ({size / (1024*1024):.1f}MB)"
                            elif size > 1024:
                                size_str = f" ({size / 1024:.0f}KB)"
                            else:
                                size_str = f" ({size}B)"
                            entries.append(("[FILE]", entry.name, size_str))
                    except OSError:
                        entries.append(("[???] ", entry.name, ""))

            if not entries:
                return f"目录为空: {path}"

            # 排序：目录在前，然后按名称
            entries.sort(key=lambda e: (e[0] != "[DIR] ", e[1].lower()))

            total = len(entries)
            if limit < total:
                entries = entries[:limit]

            lines = [f"{path} ({total} 个条目):"]
            for prefix, name, size in entries:
                lines.append(f"{prefix}{name}{size}")

            if limit < total:
                lines.append(f"... (仅显示前 {limit} 个)")

            return "\n".join(lines)

        except PermissionError:
            return f"错误: 无权限访问 - {path}"
        except Exception as e:
            return f"列出目录错误: {e}"


class GrepTool(BaseTool):
    """代码内容搜索工具 — 用正则表达式搜索文件内容"""

    name = "grep"
    is_parallel_safe = True
    description = (
        "在文件内容中搜索匹配正则表达式模式的行。"
        "支持文件类型过滤（glob 参数）、上下文行（-A/-B/-C）、大小写不敏感（-i）、"
        "以及三种输出模式：content（显示匹配行）、files_with_matches（只显示文件路径）、"
        "count（显示匹配计数）。默认只返回文件路径列表。"
    )
    parameters = {
        "pattern": ToolParameter(
            type="string",
            description="要搜索的正则表达式模式",
        ),
        "path": ToolParameter(
            type="string",
            description="搜索目录或文件路径，默认当前工作目录",
            default=".",
        ),
        "glob": ToolParameter(
            type="string",
            description="文件过滤 glob 模式，如 '*.py' 或 '*.{ts,tsx}'，默认搜索所有文本文件",
            default="**/*",
        ),
        "output_mode": ToolParameter(
            type="string",
            description="输出模式: 'content' 显示匹配行内容，'files_with_matches' 只显示文件路径，'count' 显示各文件匹配数",
            enum=["content", "files_with_matches", "count"],
            default="files_with_matches",
        ),
        "-A": ToolParameter(
            type="integer",
            description="显示匹配行之后的行数",
            default=0,
        ),
        "-B": ToolParameter(
            type="integer",
            description="显示匹配行之前的行数",
            default=0,
        ),
        "-C": ToolParameter(
            type="integer",
            description="显示匹配行前后的行数（同时设置 -A 和 -B）",
            default=0,
        ),
        "-i": ToolParameter(
            type="boolean",
            description="大小写不敏感搜索",
            default=False,
        ),
        "head_limit": ToolParameter(
            type="integer",
            description="最大输出行数，默认 250。设为 0 表示无限制",
            default=250,
        ),
    }
    required = ["pattern"]

    # 文本文件扩展名（不含二进制/媒体）
    _TEXT_EXTENSIONS = {
        ".py", ".js", ".ts", ".tsx", ".jsx", ".vue", ".svelte",
        ".html", ".htm", ".css", ".scss", ".sass", ".less",
        ".json", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ".md", ".markdown", ".rst", ".txt", ".csv", ".log",
        ".sh", ".bash", ".zsh", ".fish", ".ps1", ".bat", ".cmd",
        ".c", ".cpp", ".cc", ".cxx", ".h", ".hpp", ".hh",
        ".go", ".rs", ".java", ".kt", ".kts", ".scala",
        ".rb", ".php", ".swift", ".r", ".m", ".mm",
        ".sql", ".graphql", ".gql",
        ".env", ".gitignore", ".dockerignore", ".editorconfig",
        ".conf", ".service", ".socket", ".timer",
        ".vim", ".lua", ".el", ".ex", ".exs",
        ".tex", ".bib",
        ".Makefile", ".makefile", ".cmake",
        ".proto", ".thrift", ".avsc",
        ".lock", ".toml",
    }

    _MAX_FILE_SIZE = 2 * 1024 * 1024  # 跳过超过 2MB 的文件

    def _is_text_file(self, filepath: str) -> bool:
        """判断是否为可搜索的文本文件"""
        ext = os.path.splitext(filepath)[1].lower()
        if ext in self._TEXT_EXTENSIONS:
            return True
        # 无扩展名或未知扩展名：读取头部尝试判断
        try:
            with open(filepath, 'rb') as f:
                chunk = f.read(1024)
            # 检测 null 字节（二进制标志）
            if b'\x00' in chunk:
                return False
            # 尝试 UTF-8 解码
            chunk.decode('utf-8')
            return True
        except (UnicodeDecodeError, IOError):
            return False

    async def execute(self, pattern: str, path: str = ".", glob_pattern: str = "**/*",
                      output_mode: str = "files_with_matches", **kwargs) -> str:
        try:
            import re

            # 处理参数名冲突（glob 是 Python 内置模块名）
            after = kwargs.get("-A", 0)
            before = kwargs.get("-B", 0)
            context = kwargs.get("-C", 0)
            if context > 0:
                after = max(after, context)
                before = max(before, context)
            ignore_case = kwargs.get("-i", False)
            head_limit = kwargs.get("head_limit", 250)

            # 编译正则
            flags = re.IGNORECASE if ignore_case else 0
            try:
                regex = re.compile(pattern, flags)
            except re.error as e:
                return f"错误: 无效的正则表达式 - {e}"

            # 收集要搜索的文件
            base_path = os.path.abspath(path)
            files_to_search = []

            if os.path.isfile(base_path):
                files_to_search = [base_path]
            elif os.path.isdir(base_path):
                # 转换 glob 模式：如 *.py -> **/*.py
                if glob_pattern == "**/*":
                    search_glob = "**/*"
                elif glob_pattern.startswith("**/"):
                    search_glob = glob_pattern
                else:
                    search_glob = f"**/{glob_pattern}"

                for f in glob.iglob(os.path.join(base_path, search_glob), recursive=True):
                    if os.path.isfile(f):
                        files_to_search.append(f)
            else:
                return f"错误: 路径不存在 - {path}"

            if not files_to_search:
                return f"未找到匹配 '{glob_pattern}' 的文件（{path}）"

            # 搜索文件内容
            output_lines = []
            total_matches = 0
            files_with_matches = 0
            hit_limit = False

            for filepath in sorted(files_to_search):
                if total_matches >= 20000:  # 总匹配上限
                    hit_limit = True
                    break

                # 检查文件大小
                try:
                    fsize = os.path.getsize(filepath)
                    if fsize > self._MAX_FILE_SIZE:
                        continue
                except OSError:
                    continue

                # 检查是否为文本文件
                if not self._is_text_file(filepath):
                    continue

                # 读取并搜索
                try:
                    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                        file_lines = f.readlines()
                except (IOError, PermissionError):
                    continue

                file_matches = []
                for line_num, line in enumerate(file_lines, 1):
                    if regex.search(line):
                        file_matches.append(line_num)

                if not file_matches:
                    continue

                files_with_matches += 1
                total_matches += len(file_matches)

                rel_path = os.path.relpath(filepath, base_path) if os.path.isdir(base_path) else os.path.basename(filepath)

                if output_mode == "files_with_matches":
                    if head_limit > 0 and len(output_lines) >= head_limit:
                        hit_limit = True
                        break
                    output_lines.append(rel_path)

                elif output_mode == "count":
                    if head_limit > 0 and len(output_lines) >= head_limit:
                        hit_limit = True
                        break
                    output_lines.append(f"{rel_path}: {len(file_matches)}")

                elif output_mode == "content":
                    # 计算需要显示的行的范围（含上下文）
                    show_lines = set()
                    for ln in file_matches:
                        start = max(1, ln - before)
                        end = min(len(file_lines), ln + after)
                        for sl in range(start, end + 1):
                            show_lines.add(sl)

                    # 输出文件头
                    if head_limit > 0 and len(output_lines) >= head_limit:
                        hit_limit = True
                        break
                    output_lines.append(f"── {rel_path} ──")
                    if head_limit > 0 and len(output_lines) >= head_limit:
                        hit_limit = True
                        break

                    # 按连续行分组输出
                    sorted_lines = sorted(show_lines)
                    groups = []
                    group_start = sorted_lines[0]
                    prev = sorted_lines[0]
                    for ln in sorted_lines[1:]:
                        if ln == prev + 1:
                            prev = ln
                        else:
                            groups.append((group_start, prev))
                            group_start = ln
                            prev = ln
                    groups.append((group_start, prev))

                    for g_start, g_end in groups:
                        for ln in range(g_start, g_end + 1):
                            if head_limit > 0 and len(output_lines) >= head_limit:
                                hit_limit = True
                                break
                            marker = ":" if ln in file_matches else "-"
                            line_content = file_lines[ln - 1].rstrip('\n\r')
                            output_lines.append(f"{marker}{ln}: {line_content}")
                        if hit_limit:
                            break
                        # 组间分隔
                        if len(groups) > 1:
                            if head_limit > 0 and len(output_lines) >= head_limit:
                                hit_limit = True
                                break
                            output_lines.append("--")

                if hit_limit:
                    break

            # 构建结果
            if output_mode == "files_with_matches":
                if not output_lines:
                    return f"未找到匹配 '{pattern}' 的内容"
                result = f"搜索 '{pattern}'，找到 {files_with_matches} 个文件（共 {total_matches} 处匹配）：\n" + "\n".join(output_lines)
            elif output_mode == "count":
                if not output_lines:
                    return f"未找到匹配 '{pattern}' 的内容"
                result = f"搜索 '{pattern}'，各文件匹配数：\n" + "\n".join(output_lines)
            else:
                if not output_lines:
                    return f"未找到匹配 '{pattern}' 的内容"
                result = f"搜索 '{pattern}'，找到 {files_with_matches} 个文件（共 {total_matches} 处匹配）：\n" + "\n".join(output_lines)

            if hit_limit:
                result += f"\n... (输出已达上限 {head_limit} 行)"
            if total_matches >= 20000:
                result += "\n... (匹配数已达上限 20000)"

            return result

        except Exception as e:
            return f"搜索内容错误: {e}"


class ScreenshotTool(BaseTool):
    """截图工具 — 桌面截图 (mss) 和网页截图 (playwright)"""

    name = "screenshot"
    description = (
        "截取桌面或网页截图。"
        "桌面模式 (mode='desktop')：截取指定显示器的屏幕，可选 monitor 编号。"
        "网页模式 (mode='webpage')：对指定 URL 进行全页或视口截图，可选 viewport 尺寸和等待时间。"
        "截图保存为 PNG 格式，返回文件路径和尺寸信息。"
    )
    parameters = {
        "mode": ToolParameter(
            type="string",
            description="截图模式：'desktop' 桌面截图，'webpage' 网页截图",
            enum=["desktop", "webpage"],
        ),
        "url": ToolParameter(
            type="string",
            description="网页 URL（webpage 模式必填）",
            default="",
        ),
        "output_path": ToolParameter(
            type="string",
            description="输出文件路径（可选，默认自动生成时间戳文件名）",
            default="",
        ),
        "monitor": ToolParameter(
            type="integer",
            description="显示器编号，0=所有显示器（desktop 模式，默认 0）",
            default=0,
        ),
        "full_page": ToolParameter(
            type="boolean",
            description="是否全页截图（webpage 模式，默认 false）",
            default=False,
        ),
        "viewport_width": ToolParameter(
            type="integer",
            description="视口宽度，默认 1280（webpage 模式）",
            default=1280,
        ),
        "viewport_height": ToolParameter(
            type="integer",
            description="视口高度，默认 720（webpage 模式）",
            default=720,
        ),
        "wait_time": ToolParameter(
            type="integer",
            description="页面加载后等待时间，毫秒（webpage 模式，默认 1000）",
            default=1000,
        ),
        "timeout": ToolParameter(
            type="integer",
            description="页面加载超时，毫秒（webpage 模式，默认 30000）",
            default=30000,
        ),
    }
    required = ["mode"]

    async def execute(self, mode: str, **kwargs) -> str:
        if mode == "desktop":
            return await self._capture_desktop(**kwargs)
        elif mode == "webpage":
            return await self._capture_webpage(**kwargs)
        else:
            return json.dumps({
                "success": False,
                "error": f"无效模式 '{mode}'，必须是 'desktop' 或 'webpage'",
            })

    async def _capture_desktop(self, **kwargs) -> str:
        try:
            import mss
            import mss.tools
        except ImportError:
            return json.dumps({
                "success": False,
                "error": "mss 库未安装。安装命令: pip install mss",
            })

        monitor_num = kwargs.get("monitor", 0)
        output_path = kwargs.get("output_path", "")

        try:
            with mss.mss() as sct:
                monitors = sct.monitors
                if monitor_num < 0 or monitor_num >= len(monitors):
                    return json.dumps({
                        "success": False,
                        "error": f"无效的显示器编号 {monitor_num}，可用范围: 0-{len(monitors)-1}",
                    })

                monitor = monitors[monitor_num]
                screenshot = sct.grab(monitor)

                if not output_path:
                    from datetime import datetime
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_path = f"screenshots/desktop_{timestamp}.png"

                full_path = os.path.abspath(output_path)
                os.makedirs(os.path.dirname(full_path) or ".", exist_ok=True)
                mss.tools.to_png(screenshot.rgb, screenshot.size, output=str(full_path))

                file_size = os.path.getsize(full_path)
                return json.dumps({
                    "success": True,
                    "path": str(full_path),
                    "size": f"{screenshot.width}x{screenshot.height}",
                    "file_size": file_size,
                    "monitor": monitor_num,
                })
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})

    async def _capture_webpage(self, **kwargs) -> str:
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            return json.dumps({
                "success": False,
                "error": "playwright 库未安装。安装命令: pip install playwright && playwright install chromium",
            })

        url = kwargs.get("url", "")
        output_path = kwargs.get("output_path", "")
        full_page = kwargs.get("full_page", False)
        viewport_width = kwargs.get("viewport_width", 1280)
        viewport_height = kwargs.get("viewport_height", 720)
        wait_time = kwargs.get("wait_time", 1000)
        timeout = kwargs.get("timeout", 30000)

        if not url:
            return json.dumps({"success": False, "error": "webpage 模式必须提供 url 参数"})
        if not url.startswith(("http://", "https://")):
            return json.dumps({"success": False, "error": f"无效 URL '{url}'，必须以 http:// 或 https:// 开头"})

        try:
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                context = await browser.new_context(
                    viewport={"width": viewport_width, "height": viewport_height}
                )
                page = await context.new_page()
                await page.goto(url, wait_until="networkidle", timeout=timeout)

                if wait_time > 0:
                    await asyncio.sleep(wait_time / 1000)

                if not output_path:
                    from datetime import datetime
                    from urllib.parse import urlparse
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    domain = urlparse(url).netloc.replace(".", "_")
                    output_path = f"screenshots/webpage_{domain}_{timestamp}.png"

                full_path = os.path.abspath(output_path)
                os.makedirs(os.path.dirname(full_path) or ".", exist_ok=True)

                await page.screenshot(path=str(full_path), full_page=full_page)
                page_title = await page.title()
                await browser.close()

                file_size = os.path.getsize(full_path)
                return json.dumps({
                    "success": True,
                    "path": str(full_path),
                    "url_source": url,
                    "title": page_title,
                    "viewport": f"{viewport_width}x{viewport_height}",
                    "full_page": full_page,
                    "file_size": file_size,
                })
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


# 媒体类型扩展名映射
_IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.svg', '.ico'}
_AUDIO_EXTENSIONS = {'.mp3', '.wav', '.ogg', '.flac', '.aac', '.m4a', '.wma'}
_VIDEO_EXTENSIONS = {'.mp4', '.webm', '.avi', '.mov', '.mkv', '.flv', '.wmv'}


class DisplayMediaTool(BaseTool):
    """在聊天窗口展示媒体文件（图片、音频、视频、文档）"""

    name = "display_media"
    description = (
        "在聊天窗口向用户展示媒体文件。使用此工具展示生成或下载的图片、音频、视频或文档文件。"
        "支持自动检测文件类型（image/audio/video/file），返回文件路径和访问 URL。"
    )
    parameters = {
        "file_path": ToolParameter(
            type="string",
            description="文件的绝对路径",
        ),
        "media_type": ToolParameter(
            type="string",
            description="媒体类型：'image'、'audio'、'video' 或 'file'（不指定则自动检测）",
            enum=["image", "audio", "video", "file"],
            default="",
        ),
        "caption": ToolParameter(
            type="string",
            description="可选的标题或描述",
            default="",
        ),
    }
    required = ["file_path"]

    @staticmethod
    def _detect_media_type(file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        if ext in _IMAGE_EXTENSIONS:
            return "image"
        if ext in _AUDIO_EXTENSIONS:
            return "audio"
        if ext in _VIDEO_EXTENSIONS:
            return "video"
        return "file"

    async def execute(self, file_path: str, media_type: str = "", caption: str = "", **kwargs) -> str:
        try:
            full_path = os.path.abspath(file_path)
            if not os.path.isfile(full_path):
                return json.dumps({
                    "success": False,
                    "error": f"文件不存在: {full_path}",
                })

            if not media_type:
                media_type = self._detect_media_type(full_path)

            file_size = os.path.getsize(full_path)
            file_name = os.path.basename(full_path)
            # 构建文件访问 URL
            from urllib.parse import quote
            file_url = f"/api/files?path={quote(full_path)}"

            return json.dumps({
                "success": True,
                "path": full_path,
                "url": file_url,
                "name": file_name,
                "media_type": media_type,
                "file_size": file_size,
                "caption": caption or file_name,
            })
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class ImageTool(BaseTool):
    """图片工具 — 提取图片信息或 OCR 文字识别"""

    name = "image"
    description = (
        "读取图片信息或提取图片中的文字（OCR）。"
        "mode='ocr'：使用 OCR 识别图片中的文字，支持中文/英文等多语言。"
        "mode='info'：返回图片尺寸、格式、文件大小等元信息。"
    )
    parameters = {
        "path": ToolParameter(
            type="string",
            description="图片文件路径（支持 .png, .jpg, .jpeg, .bmp, .webp, .gif）",
        ),
        "mode": ToolParameter(
            type="string",
            description="操作模式：'ocr' 文字识别，'info' 图片信息",
            enum=["ocr", "info"],
            default="ocr",
        ),
        "language": ToolParameter(
            type="string",
            description="OCR 语言代码列表，逗号分隔。如 'ch_sim'（简体中文）、'en'（英文）、'ch_sim,en'（中英混合）。仅 mode='ocr' 有效。",
            default="ch_sim",
        ),
    }
    required = ["path"]

    _SUPPORTED_EXTENSIONS = _IMAGE_EXTENSIONS  # 复用已有的图片扩展名集合

    async def execute(self, path: str, mode: str = "ocr", language: str = "ch_sim", **kwargs) -> str:
        try:
            full_path = os.path.abspath(path)
            if not os.path.isfile(full_path):
                return json.dumps({"success": False, "error": f"文件不存在: {path}"})

            ext = os.path.splitext(full_path)[1].lower()
            if ext not in self._SUPPORTED_EXTENSIONS:
                return json.dumps({
                    "success": False,
                    "error": f"不支持的图片格式: {ext}。支持: {', '.join(sorted(self._SUPPORTED_EXTENSIONS))}"
                })

            if mode == "info":
                return await self._get_image_info(full_path)
            else:
                return await self._ocr_image(full_path, language)

        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})

    async def _get_image_info(self, path: str) -> str:
        """获取图片元信息"""
        try:
            from PIL import Image
        except ImportError:
            return json.dumps({"success": False, "error": "需要安装 Pillow 库 (pip install Pillow)"})

        try:
            with Image.open(path) as img:
                info = {
                    "success": True,
                    "path": path,
                    "format": img.format,
                    "mode": img.mode,
                    "width": img.width,
                    "height": img.height,
                    "file_size": os.path.getsize(path),
                }
                # 尝试读取 EXIF 信息（如有）
                exif = img.getexif()
                if exif:
                    exif_data = {}
                    for key, val in exif.items():
                        if isinstance(val, (str, int, float)):
                            exif_data[str(key)] = val
                    if exif_data:
                        info["exif"] = exif_data
                return json.dumps(info, ensure_ascii=False, default=str)
        except Exception as e:
            return json.dumps({"success": False, "error": f"无法读取图片: {e}"})

    async def _ocr_image(self, path: str, language: str = "ch_sim") -> str:
        """使用 OCR 提取图片中的文字"""
        try:
            import easyocr
        except ImportError:
            return json.dumps({
                "success": False,
                "error": "需要安装 easyocr 库才能进行 OCR 识别 (pip install easyocr)"
            })

        try:
            # 解析语言列表
            langs = [lang.strip() for lang in language.split(",") if lang.strip()]
            if not langs:
                langs = ["ch_sim"]

            # 初始化 Reader（懒加载：每次新建，避免占用大量内存）
            reader = easyocr.Reader(langs, gpu=False)

            # 执行 OCR
            results = reader.readtext(path)

            if not results:
                return json.dumps({
                    "success": True,
                    "mode": "ocr",
                    "language": language,
                    "text": "",
                    "blocks": [],
                    "message": "图片中未检测到文字"
                }, ensure_ascii=False)

            # 提取完整文本和按块分组
            blocks = []
            all_text = []
            for i, (bbox, text, confidence) in enumerate(results):
                blocks.append({
                    "id": i + 1,
                    "text": text,
                    "confidence": round(confidence, 4),
                    "bbox": [[int(x), int(y)] for x, y in bbox],
                })
                all_text.append(text)

            return json.dumps({
                "success": True,
                "mode": "ocr",
                "language": language,
                "text": "\n".join(all_text),
                "blocks": blocks,
                "block_count": len(blocks),
            }, ensure_ascii=False)

        except Exception as e:
            error_msg = str(e)
            if "Invalid language" in error_msg or "language" in error_msg.lower():
                error_msg = f"不支持的语言代码: {language}。常用: ch_sim (简体中文), en (英文), ch_sim,en (中英混合)"
            return json.dumps({"success": False, "error": error_msg})


class AskUserQuestionTool(BaseTool):
    """向用户提问获取决策或信息"""

    name = "ask_user_question"
    description = (
        "向用户提问以获取决策或额外信息。当你需要在多个选项中做出选择、确认操作、"
        "或需要用户提供更多上下文时使用此工具。支持单选、多选和自由文本输入。"
        "用户的选择将直接返回，用于指导后续操作。"
    )
    parameters = {
        "questions": ToolParameter(
            type="string",
            description=(
                "JSON 格式的问题列表。每个问题包含以下字段：\n"
                "- question (必填): 向用户提出的问题文本\n"
                "- header (可选): 简短标签，最多12字符\n"
                "- options (必填): 选项数组，每项含 label (显示文本) 和 description (说明)\n"
                "- multiSelect (可选): 是否允许多选，默认 false\n"
                "示例: [{\"question\":\"选择处理方式\",\"header\":\"操作\","
                "\"options\":[{\"label\":\"继续\",\"description\":\"保持现有方案\"},"
                "{\"label\":\"回滚\",\"description\":\"恢复到之前版本\"}]}]"
            ),
        ),
    }
    required = ["questions"]

    async def execute(self, questions: str, **kwargs) -> str:
        import asyncio

        # 1. 解析问题列表
        try:
            parsed = json.loads(questions)
        except json.JSONDecodeError as e:
            return json.dumps({
                "success": False,
                "error": f"questions 参数不是有效的 JSON: {e}",
            }, ensure_ascii=False)

        if not isinstance(parsed, list) or len(parsed) == 0:
            return json.dumps({
                "success": False,
                "error": "questions 必须是非空 JSON 数组",
            }, ensure_ascii=False)

        # 2. 验证每个问题的结构
        for i, q in enumerate(parsed):
            if not isinstance(q, dict):
                return json.dumps({"success": False, "error": f"问题 #{i+1} 不是有效的对象"})
            if "question" not in q:
                return json.dumps({"success": False, "error": f"问题 #{i+1} 缺少 'question' 字段"})
            if "options" not in q or not isinstance(q["options"], list) or len(q["options"]) == 0:
                return json.dumps({"success": False, "error": f"问题 #{i+1} 缺少有效的 'options' 数组"})

        # 3. 构建中断消息和选项
        from app.modules.agent.interrupt import InterruptOption, InterruptReason, get_interrupt_manager

        manager = get_interrupt_manager()

        # 构建消息文本（Markdown 格式）
        msg_lines = []
        for i, q in enumerate(parsed):
            header = q.get("header", f"问题 {i+1}")
            msg_lines.append(f"### {header}")
            msg_lines.append(q["question"])
            msg_lines.append("")
            for j, opt in enumerate(q.get("options", [])):
                label = opt.get("label", f"选项 {j+1}")
                desc = opt.get("description", "")
                msg_lines.append(f"- **{label}**{' — ' + desc if desc else ''}")
            msg_lines.append("")

        message = "\n".join(msg_lines)

        # 构建 InterruptOption 列表（使用 requires_input 让用户以 JSON 返回答案）
        options = [
            InterruptOption(
                label="提交答案",
                value="approve",
                description="请以 JSON 格式提交你的答案。格式: [{\"question_index\": 0, \"answers\": [\"选项标签\"]}, ...]",
                style="primary",
                requires_input=True,
                input_placeholder='[{"question_index": 0, "answers": ["继续"]}]',
            ),
            InterruptOption(
                label="取消",
                value="reject",
                style="default",
            ),
        ]

        # 4. 创建中断
        interrupt = await manager.create_interrupt(
            reason=InterruptReason.CUSTOM,
            message=message,
            title="Agent 请求你的决策",
            details={"questions": parsed},
            options=options,
            ttl=300.0,  # 5 分钟超时
        )

        # 5. 轮询等待用户响应
        poll_interval = 1.0
        timeout = 300.0
        elapsed = 0.0

        while elapsed < timeout:
            ip = await manager.get_interrupt(interrupt.id)
            if ip is None:
                # resolve_interrupt 会从 _pending_interrupts 中删除中断，
                # 但本地 interrupt 引用与 dict 中是同一对象，状态已更新
                if interrupt.is_resolved():
                    ip = interrupt
                else:
                    return json.dumps({
                        "success": False,
                        "error": "用户提问已过期，操作已取消",
                    }, ensure_ascii=False)
            if ip.is_expired():
                return json.dumps({
                    "success": False,
                    "error": "用户提问已过期，操作已取消",
                }, ensure_ascii=False)
            if ip.is_resolved():
                if ip.status.value == "rejected" or ip.status.value == "cancelled":
                    return json.dumps({
                        "success": False,
                        "error": "用户取消了提问",
                    }, ensure_ascii=False)
                # 解析用户答案
                try:
                    user_answers = json.loads(ip.resolution_note or "[]")
                except json.JSONDecodeError:
                    return json.dumps({
                        "success": True,
                        "answer": ip.resolution_note or "(用户未提供额外输入)",
                    }, ensure_ascii=False)
                return json.dumps({
                    "success": True,
                    "answers": user_answers,
                }, ensure_ascii=False)
            await asyncio.sleep(poll_interval)
            elapsed += poll_interval

        return json.dumps({
            "success": False,
            "error": "等待用户响应超时（5 分钟），操作已取消",
        }, ensure_ascii=False)


class RunBackgroundTool(BaseTool):
    """在后台异步运行工具，不阻塞对话"""

    name = "run_background"
    description = (
        "在后台异步执行一个耗时工具，不阻塞当前对话。立即返回 task_id。"
        "使用 check_task 工具查询任务状态和结果。"
        "适用场景：长时间运行的命令、大批量文件操作、视频/音频生成等。"
    )
    parameters = {
        "tool_name": ToolParameter(
            type="string",
            description="要在后台运行的工具名称（如 exec、web_search、screenshot 等）",
        ),
        "args": ToolParameter(
            type="object",
            description="传递给目标工具的参数字典，如 {\"command\": \"dir\"}。支持 dict 或 JSON 字符串。",
        ),
        "label": ToolParameter(
            type="string",
            description="可读的任务标签（可选，自动生成）",
            default="",
        ),
    }
    required = ["tool_name", "args"]

    async def execute(self, tool_name: str, args: str | dict, label: str = "", **kwargs) -> str:
        try:
            # 解析 args：兼容 dict 和 JSON 字符串两种格式
            if isinstance(args, dict):
                tool_args = args
            else:
                try:
                    tool_args = json.loads(args)
                    if not isinstance(tool_args, dict):
                        return json.dumps({"success": False, "error": "args 必须是 JSON 对象"})
                except json.JSONDecodeError as e:
                    return json.dumps({"success": False, "error": f"args JSON 解析失败: {e}"})

            # 从全局注册表获取工具
            global_registry = get_tool_registry()
            tool = global_registry.get_tool(tool_name)
            if not tool:
                return json.dumps({
                    "success": False,
                    "error": f"工具不存在: {tool_name}。可用工具: {global_registry.list_tools()}",
                })

            # 生成 task_id
            import uuid
            task_id = str(uuid.uuid4())[:8]
            label = label or f"后台任务: {tool_name}"

            # 存储任务状态
            _store_create_task(task_id, label, tool_name, tool_args)

            # 启动后台执行
            async def _run_background():
                try:
                    # Register inbox for inter-agent messaging
                    inbox = register_agent(agent_id=task_id, label=label)

                    from sqlalchemy import select

                    from app.api.chat import SimpleLLMProvider
                    from app.core.database import async_session_maker
                    from app.models.models import AIModelConfig
                    from app.modules.agent import AgentLoop
                    from app.modules.agent.prompts import get_tool_calling_prompt

                    # Get model config
                    async with async_session_maker() as session:
                        result = await session.execute(
                            select(AIModelConfig).where(
                                AIModelConfig.is_active
                            ).order_by(AIModelConfig.is_default.desc())
                        )
                        config = result.scalars().first()

                    if not config:
                        from app.modules.tools.task_store import update_task_status
                        update_task_status(task_id, "failed", error="No AI model config available")
                        return

                    # Build tool registry for the sub-agent
                    from app.modules.tools import ToolRegistry, register_builtin_tools
                    agent_tool_registry = ToolRegistry()
                    register_builtin_tools(agent_tool_registry)

                    # Build system prompt
                    tool_definitions = agent_tool_registry.get_definitions()
                    tools_desc = []
                    for td in tool_definitions:
                        func = td.get("function", {})
                        name = func.get("name", "unknown")
                        desc = func.get("description", "")
                        tools_desc.append(f"- {name}: {desc}")
                    system_prompt = get_tool_calling_prompt("\n".join(tools_desc))

                    provider = SimpleLLMProvider(config=config)

                    from app.core.security_client import security_client
                    sub_agent = AgentLoop(
                        provider=provider,
                        tools=agent_tool_registry,
                        model=config.model_name,
                        max_iterations=8,
                        temperature=config.temperature,
                        max_tokens=config.max_tokens,
                        permission_mode="prompt",
                        inbox_queue=inbox,
                        security_client=security_client,
                    )

                    _background_tasks[task_id]["status"] = "running"

                    # Run the agent with the tool task
                    task_message = (
                        f"Execute tool '{tool_name}' with arguments: "
                        f"{json.dumps(tool_args, ensure_ascii=False)}"
                    )
                    final_result = await sub_agent.process_direct(
                        message=task_message,
                        system_prompt=system_prompt,
                    )

                    _background_tasks[task_id]["status"] = "done"
                    _background_tasks[task_id]["result"] = final_result

                except Exception as e:
                    _background_tasks[task_id]["status"] = "failed"
                    _background_tasks[task_id]["error"] = str(e)
                finally:
                    unregister_agent(task_id)

            asyncio.create_task(_run_background())

            return json.dumps({
                "success": True,
                "task_id": task_id,
                "label": label,
                "tool_name": tool_name,
                "status": "running",
                "message": f"后台任务已启动。使用 check_task(task_id='{task_id}') 查询结果。",
            })
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class CheckTaskTool(BaseTool):
    """查询后台任务状态和结果"""

    name = "check_task"
    description = (
        "查询后台任务的状态和结果。如果省略 task_id，则列出所有运行中的后台任务。"
        "状态值: pending(等待中)、running(运行中)、done(已完成)、failed(失败)。"
    )
    parameters = {
        "task_id": ToolParameter(
            type="string",
            description="要查询的任务 ID（可选，省略则列出所有任务）",
            default="",
        ),
    }
    required = []

    async def execute(self, task_id: str = "", **kwargs) -> str:
        try:
            if task_id:
                # 查询单个任务
                task = _store_get_task(task_id)
                if not task:
                    return json.dumps({
                        "success": False,
                        "error": f"任务不存在: {task_id}",
                    })

                result_info = task.copy()
                # 截断过长结果
                if result_info.get("result") and len(result_info["result"]) > 2000:
                    result_info["result"] = result_info["result"][:2000] + "...(已截断)"
                return json.dumps(result_info, ensure_ascii=False)
            else:
                # 列出所有任务
                tasks = _store_list_tasks()
                return json.dumps({
                    "total": len(tasks),
                    "tasks": tasks,
                }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


# ============================================================
# Memory Tools — 使用 fish_memory 模块 (MemoryManage API)
# ============================================================

class MemorySaveTool(BaseTool):
    """保存记忆到文件系统（.md 文件 + MEMORY.md 索引）"""

    name = "memory_save"
    description = (
        "将内容保存为一条持久化记忆。记忆以 Markdown 文件形式存储，"
        "自动分类并更新索引。类型可选: user（用户偏好）、feedback（用户反馈）、"
        "project（项目知识）、reference（参考资料）。"
    )
    parameters = {
        "content": ToolParameter(
            type="string",
            description="记忆内容，支持 Markdown 格式",
        ),
        "type": ToolParameter(
            type="string",
            description="记忆类型: user, feedback, project, reference",
            default="user",
        ),
        "name": ToolParameter(
            type="string",
            description="记忆名称（可选，留空自动生成）",
            default="",
        ),
        "description": ToolParameter(
            type="string",
            description="记忆描述（可选，留空自动生成）",
            default="",
        ),
        "upsert": ToolParameter(
            type="boolean",
            description="如果同名或同主题记忆已存在，是否更新而非新建。默认 true",
            default=True,
        ),
    }
    required = ["content"]

    async def execute(self, content: str, type: str = "user",
                      name: str = "", description: str = "",
                      upsert: bool = True, **kwargs) -> str:
        try:
            from pathlib import Path

            from app.modules.memory import (
                MemoryFailure,
                MemoryMetadata,
                MemoryType,
                create_memory_manager,
                get_current_memory_manager,
            )

            mm = get_current_memory_manager()
            if mm is None:
                memory_root = Path(__file__).resolve().parent.parent.parent.parent / "memory"
                mm = create_memory_manager(str(memory_root))

            meta = MemoryMetadata(
                name=name or content[:50],
                description=description or content[:100],
                type=MemoryType(type),
            )

            response = mm.save(content, type, meta, upsert=upsert)
            if isinstance(response, MemoryFailure):
                err = response.error
                return json.dumps({"success": False, "error": err.message if err else "unknown"})

            entry = response.data
            return json.dumps({
                "success": True,
                "filename": entry.filename,
                "name": entry.name,
                "type": entry.type.value,
                "message": f"记忆已保存: {entry.name}",
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class MemoryRetrieveTool(BaseTool):
    """从记忆库中检索相关记忆"""

    name = "memory_retrieve"
    description = (
        "根据查询内容从记忆库中检索相关记忆。使用语义相关性排序，"
        "返回最相关的记忆内容。适用于在对话中获取上下文相关的历史信息。"
    )
    parameters = {
        "query": ToolParameter(
            type="string",
            description="检索查询，描述需要查找的信息",
        ),
        "max_results": ToolParameter(
            type="integer",
            description="最大返回条数，默认 5",
            default=5,
        ),
    }
    required = ["query"]

    async def execute(self, query: str, max_results: int = 5, **kwargs) -> str:
        try:
            from pathlib import Path

            from app.modules.memory import (
                RecallOptions,
                create_memory_manager,
                get_current_memory_manager,
            )

            mm = get_current_memory_manager()
            if mm is None:
                memory_root = Path(__file__).resolve().parent.parent.parent.parent / "memory"
                mm = create_memory_manager(str(memory_root))

            attachments = mm.recall(query, RecallOptions(max_results=max_results))
            if not attachments:
                return json.dumps({"results": [], "total": 0, "query": query})

            formatted = []
            for att in attachments:
                formatted.append({
                    "filename": att.entry.filename,
                    "name": att.entry.name,
                    "description": att.entry.description,
                    "type": att.entry.type.value,
                    "content": att.entry.content[:500],
                    "relevance": round(att.relevance_score, 3),
                    "freshness": att.entry.freshness,
                })

            return json.dumps({
                "results": formatted,
                "total": len(formatted),
                "query": query,
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class MemorySearchTool(BaseTool):
    """全文搜索记忆"""

    name = "memory_search"
    description = (
        "在记忆文件中进行全文关键词搜索。返回匹配的文件列表和内容摘要。"
        "适用于精确查找包含特定关键词的记忆。"
    )
    parameters = {
        "keyword": ToolParameter(
            type="string",
            description="搜索关键词",
        ),
        "limit": ToolParameter(
            type="integer",
            description="最大返回条数，默认 10",
            default=10,
        ),
    }
    required = ["keyword"]

    async def execute(self, keyword: str, limit: int = 10, **kwargs) -> str:
        try:
            from pathlib import Path

            from app.modules.memory import (
                SearchOptions,
                create_memory_manager,
                get_current_memory_manager,
            )

            mm = get_current_memory_manager()
            if mm is None:
                memory_root = Path(__file__).resolve().parent.parent.parent.parent / "memory"
                mm = create_memory_manager(str(memory_root))

            response = mm.search(keyword, SearchOptions(limit=limit))
            entries = response.data if hasattr(response, 'data') else []

            formatted = []
            for e in entries:
                formatted.append({
                    "filename": e.filename,
                    "name": e.name,
                    "description": e.description,
                    "type": e.type.value,
                    "content_preview": e.content[:300],
                })

            return json.dumps({
                "results": formatted,
                "total": len(formatted),
                "keyword": keyword,
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class MemoryListTool(BaseTool):
    """列出所有记忆"""

    name = "memory_list"
    description = (
        "列出记忆库中的所有记忆条目，可按类型过滤。返回名称、描述、类型、新鲜度等信息。"
    )
    parameters = {
        "type": ToolParameter(
            type="string",
            description="按类型过滤: user, feedback, project, reference（可选）",
            default="",
        ),
    }
    required = []

    async def execute(self, type: str = "", **kwargs) -> str:
        try:
            from pathlib import Path

            from app.modules.memory import (
                ListOptions,
                MemoryType,
                create_memory_manager,
                get_current_memory_manager,
            )

            mm = get_current_memory_manager()
            if mm is None:
                memory_root = Path(__file__).resolve().parent.parent.parent.parent / "memory"
                mm = create_memory_manager(str(memory_root))

            mem_type = MemoryType(type) if type else None
            response = mm.list(ListOptions(type=mem_type, limit=50))
            entries = response.data if hasattr(response, 'data') else []

            formatted = []
            for e in entries:
                formatted.append({
                    "filename": e.filename,
                    "name": e.name,
                    "description": e.description,
                    "type": e.type.value,
                    "freshness": e.freshness,
                    "is_stale": e.is_stale,
                    "tags": e.tags,
                })

            return json.dumps({
                "results": formatted,
                "total": len(formatted),
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class SpawnTool(BaseTool):
    """启动子 Agent 处理复杂任务"""

    name = "spawn"
    description = (
        "启动一个子 Agent 来处理复杂的多步骤任务。子 Agent 拥有独立的工具调用能力，"
        "可以自主搜索、读写文件、执行命令等。完成后返回结果。"
        "适用场景：需要多步推理和工具调用的复杂任务、独立的研究/分析任务。"
        "超时时间 120 秒。"
        "\n\n内置 Agent 类型：\n"
        "- general（默认）: 完整工具集，适合大多数任务\n"
        "- explore: 只读探索工具（read_file/grep/file_search/list_dir/web_search），适合代码库探索\n"
        "- plan: 只读工具 + 计划模式（enter/exit_plan_mode），适合设计方案\n"
        "- verification: 只读工具 + exec，适合运行测试和验证代码变更\n"
        "- guide: 只读工具 + 知识库查询，适合回答 PioneClaw 使用问题"
    )
    parameters = {
        "task": ToolParameter(
            type="string",
            description="分配给子 Agent 的任务描述，越详细越好",
        ),
        "context": ToolParameter(
            type="string",
            description="子 Agent 的背景信息或上下文（可选）",
            default="",
        ),
        "agent_type": ToolParameter(
            type="string",
            description=(
                "子 Agent 类型：'general'（完整工具集）、'explore'（只读探索）、'plan'（计划模式）、"
                "'verification'（测试验证）、'guide'（使用指南）。默认 'general'"
            ),
            enum=["general", "explore", "plan", "verification", "guide"],
            default="general",
        ),
    }
    required = ["task"]

    # 各 Agent 类型的工具白名单（None = 全部工具）
    _AGENT_TYPE_TOOLS: dict = {
        "general": None,
        "explore": {
            "read_file", "grep", "file_search", "list_dir",
            "web_search", "web_fetch", "current_time",
        },
        "plan": {
            "read_file", "grep", "file_search", "list_dir",
            "web_search", "web_fetch", "current_time", "calculator",
            "ask_user_question", "enter_plan_mode", "exit_plan_mode",
            "image", "memory_retrieve", "memory_search",
        },
        "verification": {
            "read_file", "grep", "file_search", "list_dir",
            "web_search", "web_fetch", "current_time",
            "exec",
        },
        "guide": {
            "read_file", "grep", "file_search", "list_dir",
            "web_search", "web_fetch", "current_time",
            "memory_retrieve", "memory_search",
        },
    }

    async def execute(self, task: str, context: str = "", **kwargs) -> str:
        try:
            from sqlalchemy import select

            from app.core.database import async_session_maker
            from app.models.models import AIModelConfig

            # 获取默认模型配置
            async with async_session_maker() as session:
                result = await session.execute(
                    select(AIModelConfig).where(
                        AIModelConfig.is_active
                    ).order_by(AIModelConfig.is_default.desc())
                )
                config = result.scalars().first()

            if not config:
                return json.dumps({"success": False, "error": "没有可用的 AI 模型配置"})

            # 构建子 Agent
            from app.modules.agent import AgentLoop
            from app.modules.agent.prompts import get_tool_calling_prompt
            from app.modules.tools import ToolRegistry, register_builtin_tools

            tool_registry = ToolRegistry()
            register_builtin_tools(tool_registry)
            tool_definitions = tool_registry.get_definitions()

            # 按 Agent 类型过滤工具
            agent_type = kwargs.get("agent_type", "general")
            tool_filter = self._AGENT_TYPE_TOOLS.get(agent_type)
            if tool_filter is not None:
                tool_definitions = [
                    t for t in tool_definitions
                    if t.get("function", {}).get("name") in tool_filter
                ]

            tools_desc = []
            for td in tool_definitions:
                func = td.get("function", {})
                name = func.get("name", "unknown")
                desc = func.get("description", "")
                tools_desc.append(f"- {name}: {desc}")
            system_prompt = get_tool_calling_prompt("\n".join(tools_desc))

            # 按 Agent 类型注入专用指令
            if agent_type == "explore":
                system_prompt += (
                    "\n\n你是一个**代码库探索助手**。你的职责是搜索、阅读和理解代码，"
                    "而不是修改代码。请提供清晰、结构化的发现报告，"
                    "包括相关文件路径、关键代码片段和你的分析结论。"
                )
            elif agent_type == "plan":
                system_prompt += (
                    "\n\n你是一个**方案设计助手**。使用 enter_plan_mode 进入计划模式"
                    "进行只读调研，设计方案后使用 exit_plan_mode 提交计划。"
                    "计划应包含：Context（背景）、改动方案（具体步骤）、涉及文件、验证方法。"
                )
            elif agent_type == "verification":
                system_prompt += (
                    "\n\n你是一个**测试验证助手**。你的职责是运行测试、验证代码变更的正确性。"
                    "请运行相关的单元测试、集成测试或 E2E 测试，"
                    "分析测试结果，报告通过/失败情况，"
                    "并对失败原因给出初步诊断。"
                )
            elif agent_type == "guide":
                system_prompt += (
                    "\n\n你是一个**PioneClaw 使用指南助手**。你的职责是回答关于 PioneClaw 的使用问题。"
                    "请搜索代码库和文档，找到相关配置、功能说明或最佳实践，"
                    "给出清晰、实用的指导。"
                )

            from app.api.chat import SimpleLLMProvider
            provider = SimpleLLMProvider(config=config)

            from app.core.security_client import security_client
            sub_agent = AgentLoop(
                provider=provider,
                tools=tool_registry,
                model=config.model_name,
                max_iterations=8,
                temperature=config.temperature,
                max_tokens=config.max_tokens,
                permission_mode="prompt",  # 子 Agent 使用 Prompt 模式，自动审批
                security_client=security_client,
            )

            full_task = task
            if context:
                full_task = f"背景信息:\n{context}\n\n任务:\n{task}"

            try:
                result = await asyncio.wait_for(
                    sub_agent.process_direct(
                        message=full_task,
                        system_prompt=system_prompt,
                    ),
                    timeout=120.0,
                )

                return json.dumps({
                    "success": True,
                    "result": result,
                    "iterations": sub_agent.max_iterations,
                }, ensure_ascii=False)
            except asyncio.TimeoutError:
                return json.dumps({
                    "success": False,
                    "error": "子 Agent 执行超时（120 秒）",
                })
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class CronTool(BaseTool):
    """管理定时任务"""

    name = "cron"
    description = (
        "管理定时任务（cron jobs）。支持查看任务列表、查看任务详情、验证 cron 表达式、"
        "查看调度器状态。操作类型: list(列表)、get(详情)、validate(验证表达式)、status(调度器状态)。"
    )
    parameters = {
        "action": ToolParameter(
            type="string",
            description="操作类型: list(列出所有任务)、get(查看单个任务)、validate(验证cron表达式)、status(调度器状态)",
            enum=["list", "get", "validate", "status"],
        ),
        "job_id": ToolParameter(
            type="string",
            description="任务 ID（get 操作必填）",
            default="",
        ),
        "cron_expr": ToolParameter(
            type="string",
            description="cron 表达式（validate 操作必填），如 '0 9 * * *' 表示每天9点",
            default="",
        ),
    }
    required = ["action"]

    async def execute(self, action: str, job_id: str = "", cron_expr: str = "", **kwargs) -> str:
        try:
            from app.core.cron_scheduler import get_cron_scheduler

            scheduler = get_cron_scheduler()

            if action == "list":
                jobs = scheduler.list_jobs()
                formatted = []
                for j in jobs:
                    formatted.append({
                        "job_id": j.get("job_id", ""),
                        "cron_expr": j.get("cron_expr", ""),
                        "enabled": j.get("enabled", False),
                        "next_run": j.get("next_run", ""),
                        "last_run": j.get("last_run", ""),
                        "run_count": j.get("run_count", 0),
                    })
                return json.dumps({"jobs": formatted, "total": len(formatted)}, ensure_ascii=False)

            elif action == "get":
                if not job_id:
                    return json.dumps({"success": False, "error": "get 操作需要提供 job_id"})
                job = scheduler.get_job(job_id)
                if not job:
                    return json.dumps({"success": False, "error": f"任务不存在: {job_id}"})
                return json.dumps(job, ensure_ascii=False, default=str)

            elif action == "validate":
                if not cron_expr:
                    return json.dumps({"success": False, "error": "validate 操作需要提供 cron_expr"})
                valid = scheduler.validate_cron_expr(cron_expr)
                result = {"valid": valid, "cron_expr": cron_expr}
                if valid:
                    result["description"] = scheduler.describe_cron_expr(cron_expr)
                    next_run = scheduler.get_next_run(cron_expr)
                    result["next_run"] = str(next_run) if next_run else None
                return json.dumps(result, ensure_ascii=False)

            elif action == "status":
                jobs = scheduler.list_jobs()
                enabled = sum(1 for j in jobs if j.get("enabled"))
                return json.dumps({
                    "total_jobs": len(jobs),
                    "enabled_jobs": enabled,
                    "disabled_jobs": len(jobs) - enabled,
                }, ensure_ascii=False)

            else:
                return json.dumps({"success": False, "error": f"未知操作: {action}"})
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class CronCreateTool(BaseTool):
    """创建定时任务，根据 cron 表达式定时执行指定的提示词/任务"""

    name = "cron_create"
    description = (
        "创建定时任务。根据 cron 表达式定时执行指定的提示词/任务。"
        "cron 表达式为标准 5 字段格式：分 时 日 月 周，例如 '0 9 * * *' 表示每天 9 点。"
    )
    parameters = {
        "cron_expr": ToolParameter(
            type="string",
            description="5 字段 cron 表达式，如 '0 9 * * *' 表示每天 9 点",
        ),
        "prompt": ToolParameter(
            type="string",
            description="定时执行的提示词/任务描述",
        ),
        "description": ToolParameter(
            type="string",
            description="任务描述",
            default="",
        ),
        "name": ToolParameter(
            type="string",
            description="任务名称（可选，不填则自动生成）",
            default="",
        ),
    }
    required = ["cron_expr", "prompt"]

    async def execute(self, cron_expr: str, prompt: str, description: str = "",
                      name: str = "", **kwargs) -> str:
        try:
            from app.core.cron_scheduler import get_cron_scheduler
            from app.core.database import async_session_maker
            from app.models.models import CronJob

            scheduler = get_cron_scheduler()

            # 验证 cron 表达式
            if not scheduler.validate_cron_expr(cron_expr):
                return json.dumps({
                    "success": False,
                    "error": f"无效的 cron 表达式: {cron_expr}"
                })

            # 生成任务名称
            job_name = name.strip() if name.strip() else f"cron_{uuid.uuid4().hex[:8]}"

            # 添加到内存调度器（使用真实 Agent 回调）
            from app.core.cron_scheduler import _make_cron_callback
            callback = _make_cron_callback(job_name, {"prompt": prompt})

            success = scheduler.add_job(job_name, cron_expr, callback, enabled=True)
            if not success:
                return json.dumps({"success": False, "error": "添加到调度器失败"})

            job_info = scheduler.get_job(job_name)

            # 持久化到数据库
            db_job_id = None
            try:
                async with async_session_maker() as session:
                    db_job = CronJob(
                        name=job_name,
                        display_name=job_name,
                        schedule_type="cron",
                        schedule_value=cron_expr,
                        description=description or prompt,
                        is_enabled=True,
                        config={"prompt": prompt},
                    )
                    session.add(db_job)
                    await session.commit()
                    await session.refresh(db_job)
                    db_job_id = db_job.id
            except Exception as e:
                logger.warning(f"[CronCreateTool] DB 持久化失败（调度器已添加）: {e}")

            return json.dumps({
                "success": True,
                "job_id": job_name,
                "db_id": db_job_id,
                "cron_expr": cron_expr,
                "prompt": prompt,
                "description": description or prompt,
                "next_run": str(job_info["next_run"]) if job_info and job_info.get("next_run") else None,
                "message": f"定时任务 '{job_name}' 创建成功"
            }, ensure_ascii=False)

        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class CronDeleteTool(BaseTool):
    """删除/取消定时任务"""

    name = "cron_delete"
    description = (
        "删除/取消定时任务。从调度器和数据库中移除指定的定时任务。"
    )
    parameters = {
        "job_id": ToolParameter(
            type="string",
            description="要删除的任务 ID（名称）",
        ),
    }
    required = ["job_id"]

    async def execute(self, job_id: str, **kwargs) -> str:
        try:
            from sqlalchemy import select

            from app.core.cron_scheduler import get_cron_scheduler
            from app.core.database import async_session_maker
            from app.models.models import CronJob

            scheduler = get_cron_scheduler()

            # 从内存调度器移除
            removed = scheduler.remove_job(job_id)

            # 从数据库删除
            db_deleted = False
            try:
                async with async_session_maker() as session:
                    result = await session.execute(
                        select(CronJob).where(CronJob.name == job_id)
                    )
                    db_job = result.scalar_one_or_none()
                    if db_job:
                        await session.delete(db_job)
                        await session.commit()
                        db_deleted = True
            except Exception as e:
                logger.warning(f"[CronDeleteTool] DB 删除失败: {e}")

            if not removed and not db_deleted:
                return json.dumps({
                    "success": False,
                    "error": f"任务不存在: {job_id}"
                })

            return json.dumps({
                "success": True,
                "job_id": job_id,
                "scheduler_removed": removed,
                "db_deleted": db_deleted,
                "message": f"定时任务 '{job_id}' 已删除"
            }, ensure_ascii=False)

        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class CronListTool(BaseTool):
    """列出所有定时任务及其状态"""

    name = "cron_list"
    description = (
        "列出所有定时任务及其状态。显示任务 ID、cron 表达式、启用状态、下次执行时间等信息。"
    )
    parameters = {
        "enabled_only": ToolParameter(
            type="boolean",
            description="仅显示启用中的任务",
            default=False,
        ),
    }
    required = []

    async def execute(self, enabled_only: bool = False, **kwargs) -> str:
        try:
            from app.core.cron_scheduler import get_cron_scheduler

            scheduler = get_cron_scheduler()
            jobs = scheduler.list_jobs()

            # 过滤
            if enabled_only:
                jobs = [j for j in jobs if j.get("enabled")]

            formatted = []
            for j in jobs:
                formatted.append({
                    "job_id": j.get("job_id", ""),
                    "cron_expr": j.get("cron_expr", ""),
                    "enabled": j.get("enabled", False),
                    "next_run": str(j.get("next_run", "")) if j.get("next_run") else "",
                    "last_run": str(j.get("last_run", "")) if j.get("last_run") else "",
                    "run_count": j.get("run_count", 0),
                })

            return json.dumps({
                "jobs": formatted,
                "total": len(formatted),
                "filtered": enabled_only,
            }, ensure_ascii=False)

        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


class ChannelTool(BaseTool):
    """通过消息通道发送消息"""

    name = "channel"
    description = (
        "通过已配置的消息通道发送消息。支持飞书、钉钉、企业微信、QQ、微信等。"
        "操作类型: list(列出已配置通道)、send(发送消息)、broadcast(广播到所有通道)、status(通道状态)。"
    )
    parameters = {
        "action": ToolParameter(
            type="string",
            description="操作类型: list(列出通道)、send(发送消息)、broadcast(广播)、status(通道状态)",
            enum=["list", "send", "broadcast", "status"],
        ),
        "channel_id": ToolParameter(
            type="string",
            description="通道 ID（send 操作必填）",
            default="",
        ),
        "chat_id": ToolParameter(
            type="string",
            description="目标聊天 ID（send 操作必填）",
            default="",
        ),
        "content": ToolParameter(
            type="string",
            description="要发送的消息内容（send/broadcast 操作必填）",
            default="",
        ),
    }
    required = ["action"]

    async def execute(self, action: str, channel_id: str = "", chat_id: str = "",
                      content: str = "", **kwargs) -> str:
        try:
            from app.modules.channels import get_channel_manager

            manager = get_channel_manager()

            if action == "list":
                ids = manager.get_channel_ids()
                connected = set(manager.get_connected_channels())
                channels = []
                for cid in ids:
                    channels.append({
                        "channel_id": cid,
                        "connected": cid in connected,
                    })
                return json.dumps({"channels": channels, "total": len(channels)}, ensure_ascii=False)

            elif action == "send":
                if not channel_id or not chat_id or not content:
                    return json.dumps({
                        "success": False,
                        "error": "send 操作需要提供 channel_id、chat_id 和 content",
                    })
                success, msg_id = await manager.send_message(channel_id, chat_id, content)
                return json.dumps({
                    "success": success,
                    "message_id": msg_id,
                    "channel_id": channel_id,
                }, ensure_ascii=False)

            elif action == "broadcast":
                if not content:
                    return json.dumps({
                        "success": False,
                        "error": "broadcast 操作需要提供 content",
                    })
                results = await manager.broadcast(content)
                formatted = {}
                for cid, (success, msg_id) in results.items():
                    formatted[cid] = {"success": success, "message_id": msg_id}
                return json.dumps({"results": formatted}, ensure_ascii=False)

            elif action == "status":
                status = manager.get_status()
                return json.dumps(status, ensure_ascii=False, default=str)

            else:
                return json.dumps({"success": False, "error": f"未知操作: {action}"})
        except Exception as e:
            return json.dumps({"success": False, "error": str(e)})


# 浏览器页面缓存（模块级别，同一会话内复用）
_browser_page = None
_browser_instance = None


class BrowserTool(BaseTool):
    """浏览器自动化工具 — 打开网页、提取内容、点击交互"""

    name = "browser"
    description = (
        "浏览器自动化工具，用于打开网页、提取内容和交互操作。\n"
        "支持的操作：\n"
        "- content: 打开 URL 并提取页面文本（清洗后的正文）\n"
        "- links: 打开 URL 并提取所有链接\n"
        "- extract: 打开 URL，按 CSS 选择器提取指定元素内容\n"
        "- click: 点击页面中匹配文本的元素，返回点击后的页面内容\n"
        "- fill: 在输入框中填入文本（需先 navigate 打开页面）\n"
        "- screenshot: 对当前页面截图\n"
        "- close: 关闭浏览器"
    )
    parameters = {
        "action": ToolParameter(
            type="string",
            description="操作类型: content、links、extract、click、fill、screenshot、close",
            enum=["content", "links", "extract", "click", "fill", "screenshot", "close"],
        ),
        "url": ToolParameter(
            type="string",
            description="目标 URL（content/links/extract 操作需要）",
            default="",
        ),
        "selector": ToolParameter(
            type="string",
            description="CSS 选择器（extract 操作需要），如 'div.article'、'p'",
            default="",
        ),
        "target_text": ToolParameter(
            type="string",
            description="要点击或填写的目标文本（click: 按钮/链接文本; fill: 要填入的内容）",
            default="",
        ),
        "input_selector": ToolParameter(
            type="string",
            description="输入框的 CSS 选择器（fill 操作需要），如 'input[name=q]'、'#search'",
            default="",
        ),
        "output_path": ToolParameter(
            type="string",
            description="截图输出路径（screenshot 操作可选）",
            default="",
        ),
    }
    required = ["action"]

    async def execute(self, action: str, url: str = "", selector: str = "",
                      target_text: str = "", input_selector: str = "",
                      output_path: str = "", **kwargs) -> str:
        global _browser_page, _browser_instance

        try:
            from playwright.async_api import async_playwright
        except ImportError:
            return json.dumps({
                "success": False,
                "error": "playwright 未安装。安装命令: pip install playwright && playwright install chromium",
            })

        try:
            # 关闭浏览器
            if action == "close":
                if _browser_instance:
                    await _browser_instance.close()
                    _browser_instance = None
                    _browser_page = None
                return json.dumps({"success": True, "message": "浏览器已关闭"})

            # 需要页面的操作：确保浏览器已启动
            if action in ("content", "links", "extract", "click", "fill", "screenshot"):
                if not url and action in ("content", "links", "extract"):
                    return json.dumps({"success": False, "error": f"{action} 操作需要提供 url"})

                # 启动或复用浏览器
                if _browser_instance is None:
                    _browser_instance = await async_playwright().__aenter__()
                    browser = await _browser_instance.chromium.launch(headless=True)
                    context = await browser.new_context(
                        viewport={"width": 1280, "height": 720},
                        user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
                    )
                    _browser_page = await context.new_page()
                else:
                    browser = _browser_instance

                page = _browser_page

                # navigate 如果需要
                if url and action in ("content", "links", "extract"):
                    await page.goto(url, wait_until="networkidle", timeout=30000)
                    await asyncio.sleep(0.5)

                if action == "content":
                    # 提取页面文本（从 body 中获取）
                    text = await page.evaluate("""() => {
                        const body = document.body;
                        if (!body) return '';
                        const clone = body.cloneNode(true);
                        // 移除 script, style, nav, footer, header
                        clone.querySelectorAll('script,style,nav,footer,header,aside,noscript,iframe').forEach(e => e.remove());
                        return clone.innerText;
                    }""")
                    title = await page.title()
                    text = re.sub(r'\n{3,}', '\n\n', text)  # 压缩空行
                    if len(text) > 5000:
                        text = text[:5000] + "\n... (内容已截断)"
                    return json.dumps({
                        "success": True,
                        "url": page.url,
                        "title": title,
                        "content": text,
                    }, ensure_ascii=False)

                elif action == "links":
                    links = await page.evaluate("""() => {
                        return Array.from(document.querySelectorAll('a[href]')).map(a => ({
                            text: a.innerText.trim().substring(0, 100),
                            href: a.href
                        })).filter(l => l.href && !l.href.startsWith('javascript:'));
                    }""")
                    if len(links) > 100:
                        links = links[:100]
                    return json.dumps({
                        "success": True,
                        "url": page.url,
                        "title": await page.title(),
                        "links": links,
                        "total": len(links),
                    }, ensure_ascii=False)

                elif action == "extract":
                    if not selector:
                        return json.dumps({"success": False, "error": "extract 操作需要提供 selector"})
                    elements = await page.evaluate("""(sel) => {
                        return Array.from(document.querySelectorAll(sel)).map(el => ({
                            text: el.innerText.trim().substring(0, 500),
                            tag: el.tagName.toLowerCase()
                        }));
                    }""", selector)
                    if len(elements) > 50:
                        elements = elements[:50]
                    return json.dumps({
                        "success": True,
                        "selector": selector,
                        "elements": elements,
                        "total": len(elements),
                    }, ensure_ascii=False)

                elif action == "click":
                    if not target_text:
                        return json.dumps({"success": False, "error": "click 操作需要提供 target_text"})
                    # 点击包含指定文本的元素
                    clicked = await page.evaluate("""(text) => {
                        const elements = document.querySelectorAll('a,button,input[type=submit],input[type=button],[role=button]');
                        for (const el of elements) {
                            if (el.innerText && el.innerText.trim().includes(text)) {
                                el.click();
                                return el.innerText.trim().substring(0, 200);
                            }
                        }
                        return null;
                    }""", target_text)
                    if not clicked:
                        return json.dumps({"success": False, "error": f"未找到包含 '{target_text}' 的可点击元素"})
                    await asyncio.sleep(1)
                    text = await page.evaluate("() => document.body ? document.body.innerText.substring(0, 3000) : ''")
                    return json.dumps({
                        "success": True,
                        "clicked": clicked,
                        "url": page.url,
                        "content": text,
                    }, ensure_ascii=False)

                elif action == "fill":
                    if not input_selector or not target_text:
                        return json.dumps({"success": False, "error": "fill 操作需要提供 input_selector 和 target_text"})
                    await page.fill(input_selector, target_text)
                    return json.dumps({
                        "success": True,
                        "selector": input_selector,
                        "filled": target_text,
                    }, ensure_ascii=False)

                elif action == "screenshot":
                    if not output_path:
                        from datetime import datetime
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        output_path = f"screenshots/browser_{timestamp}.png"
                    full_path = os.path.abspath(output_path)
                    os.makedirs(os.path.dirname(full_path) or ".", exist_ok=True)
                    await page.screenshot(path=str(full_path), full_page=False)
                    file_size = os.path.getsize(full_path)
                    return json.dumps({
                        "success": True,
                        "path": str(full_path),
                        "file_size": file_size,
                    })

            return json.dumps({"success": False, "error": f"未知操作: {action}"})

        except Exception as e:
            # 发生错误时清理浏览器状态
            error_msg = str(e)
            if "net::ERR_NAME_NOT_RESOLVED" in error_msg:
                error_msg = f"无法解析域名，请检查 URL: {url}"
            elif "net::ERR_CONNECTION_REFUSED" in error_msg:
                error_msg = f"连接被拒绝: {url}"
            elif "Timeout" in error_msg:
                error_msg = f"页面加载超时: {url}"
            return json.dumps({"success": False, "error": error_msg})


# 注册所有内置工具
def register_builtin_tools(registry: Optional["ToolRegistry"] = None):
    """注册所有内置工具

    Args:
        registry: 工具注册表，默认使用全局注册表
    """
    if registry is None:
        registry = get_tool_registry()

    registry.register_class(CurrentTimeTool)
    registry.register_class(CalculatorTool)
    registry.register_class(WebSearchTool)
    registry.register_class(ReadFileTool)
    registry.register_class(WriteFileTool)
    registry.register_class(ExecTool)
    registry.register_class(EditFileTool)
    registry.register_class(FileSearchTool)
    registry.register_class(ListDirTool)
    registry.register_class(GrepTool)
    registry.register_class(ScreenshotTool)
    registry.register_class(DisplayMediaTool)
    registry.register_class(ImageTool)
    registry.register_class(AskUserQuestionTool)
    registry.register_class(EnterPlanModeTool)
    registry.register_class(ExitPlanModeTool)
    registry.register_class(RunBackgroundTool)
    registry.register_class(CheckTaskTool)
    # Memory tools (fish_memory system)
    registry.register_class(MemorySaveTool)
    registry.register_class(MemoryRetrieveTool)
    registry.register_class(MemorySearchTool)
    registry.register_class(MemoryListTool)
    registry.register_class(SpawnTool)
    registry.register_class(CronTool)
    registry.register_class(CronCreateTool)
    registry.register_class(CronDeleteTool)
    registry.register_class(CronListTool)
    registry.register_class(MCPTool)
    registry.register_class(ListMcpResourcesTool)
    registry.register_class(ReadMcpResourceTool)
    registry.register_class(McpAuthTool)
    registry.register_class(BrowserTool)
    registry.register_class(ThinkingTool)
    registry.register_class(ViewTool)
    # 注册 web 工具
    registry.register(WebSearchTool())
    registry.register(WebFetchTool())

    # 同步注册到全局注册表（供 run_background / check_task 等工具使用）
    global_registry = get_tool_registry()
    if registry is not global_registry:
        for name in registry.list_tools():
            if name not in global_registry.list_tools():
                global_registry.register(registry.get_tool(name))

    logger.info(f"Built-in tools registered: {registry.list_tools()}")
